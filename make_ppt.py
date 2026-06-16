"""MSDS Label Maker v2.1.0 사용 설명서 PPT — 캡처 이미지 레이아웃 버전"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ── 색상 팔레트 ─────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x2B, 0x4B)
TEAL   = RGBColor(0x0D, 0x8E, 0x8E)
TEAL2  = RGBColor(0x02, 0xB0, 0xAA)
LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x6C, 0x75, 0x7D)
LGRAY  = RGBColor(0xCE, 0xD4, 0xDA)
DARK   = RGBColor(0x21, 0x25, 0x29)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
RED    = RGBColor(0xC0, 0x39, 0x2B)

# ── 헬퍼 함수 ───────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill, line_rgb=None, line_w=0.75):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_rgb:
        sh.line.color.rgb = line_rgb; sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    return sh

def circle(slide, x, y, d, fill):
    sh = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()
    return sh

def txt(slide, text, x, y, w, h, size=12, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Malgun Gothic"
    return tb

def num_circle(slide, n, x, y, d=0.42, bg=None):
    bg = bg or TEAL
    sh = circle(slide, x, y, d, bg)
    p = sh.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = "Malgun Gothic"

def header(slide, num, title, subtitle):
    rect(slide, 0, 0, 10, 0.9, NAVY)
    txt(slide, f"{num}  {title}", 0.4, 0.14, 9, 0.48, size=24, bold=True, color=WHITE)
    txt(slide, subtitle, 0.4, 0.64, 9, 0.24, size=10.5, color=TEAL2)

def screenshot_box(slide, x, y, w, h, label="캡처 이미지"):
    """점선 테두리 캡처 이미지 플레이스홀더"""
    # 배경
    bg = rect(slide, x, y, w, h, RGBColor(0xF8, 0xF9, 0xFA),
              line_rgb=LGRAY, line_w=1.0)
    # 점선 효과를 위한 XML 패치
    sp_pr = bg._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
    if sp_pr is not None:
        prstDash = etree.SubElement(sp_pr, '{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash')
        prstDash.set('val', 'dashDot')
    # 중앙 아이콘 텍스트
    cy = y + h / 2 - 0.32
    txt(slide, "📷", x, cy, w, 0.5, size=22, align=PP_ALIGN.CENTER, color=LGRAY)
    txt(slide, f"[ {label} ]", x, cy + 0.42, w, 0.28,
        size=10, color=LGRAY, align=PP_ALIGN.CENTER, italic=True)

def step_row(slide, n, title, desc, x, y, w=4.3, title_color=NAVY):
    num_circle(slide, n, x, y + 0.18)
    txt(slide, title, x + 0.52, y + 0.14, w - 0.55, 0.3,
        size=12, bold=True, color=title_color)
    txt(slide, desc, x + 0.52, y + 0.44, w - 0.55, 0.52,
        size=10, color=GRAY, wrap=True)

def bullet_box(slide, x, y, w, h, head, icon, items, hbg=NAVY):
    rect(slide, x, y, w, h, LIGHT, line_rgb=LGRAY)
    rect(slide, x, y, w, 0.38, hbg)
    txt(slide, f"{icon}  {head}", x + 0.14, y + 0.07, w - 0.2, 0.26,
        size=11, bold=True, color=WHITE)
    for i, item in enumerate(items):
        txt(slide, f"▸  {item}", x + 0.16, y + 0.5 + i * 0.38, w - 0.24, 0.34,
            size=10.5, color=DARK, wrap=True)

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
blank = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — 표지
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 10, 5.625, NAVY)
rect(s, 0, 0, 0.18, 5.625, TEAL)   # 왼쪽 틸 바
rect(s, 1.3, 0.8, 7.4, 4.0, WHITE) # 카드
rect(s, 1.3, 0.8, 7.4, 0.07, TEAL2) # 카드 상단 강조

txt(s, "🧪", 4.55, 1.1, 1.0, 0.85, size=40, align=PP_ALIGN.CENTER, color=TEAL)
txt(s, "MSDS Label Maker", 1.3, 1.9, 7.4, 0.65,
    size=32, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
txt(s, "v2.1.0", 1.3, 2.52, 7.4, 0.4,
    size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
rect(s, 3.5, 3.0, 3.0, 0.04, TEAL2)
txt(s, "사용 설명서", 1.3, 3.1, 7.4, 0.44,
    size=20, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "LX글라스 연구기획팀  |  박재영  |  2026-06-11",
    1.3, 3.7, 7.4, 0.33, size=11, color=GRAY, align=PP_ALIGN.CENTER)
rect(s, 7.9, 4.85, 1.9, 0.42, TEAL)
txt(s, "GHS 경고표지 자동화", 7.92, 4.9, 1.85, 0.3,
    size=9, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — 목차
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 10, 5.625, LIGHT)
rect(s, 0, 0, 10, 0.9, NAVY)
txt(s, "목  차", 0.4, 0.14, 6, 0.52, size=24, bold=True, color=WHITE)
txt(s, "MSDS Label Maker v2.1.0", 0.4, 0.64, 6, 0.24, size=10.5, color=TEAL2)

toc = [
    ("01", "프로그램 소개",       "도구의 목적과 핵심 기능"),
    ("02", "화면 구성",           "3패널 UI 레이아웃 안내"),
    ("03", "Step 1  PDF 불러오기","PDF 파일 선택 및 텍스트 추출"),
    ("04", "Step 2  내용 편집",   "추출 결과 확인 및 항목 수정"),
    ("05", "Step 3  인쇄 출력",   "미리보기 확인 및 인쇄 HTML 생성"),
    ("06", "라벨 규격 선택",      "아이라벨 3종 규격 안내"),
    ("07", "GHS 그림문자",        "9종 그림문자 선택 방법"),
    ("08", "주의사항 & FAQ",      "올바른 사용을 위한 체크리스트"),
]
for i, (num, title, desc) in enumerate(toc):
    cx = 0.4 + (i % 2) * 4.8
    ry = 1.1 + (i // 2) * 1.1
    rect(s, cx, ry, 4.35, 0.92, WHITE, line_rgb=LGRAY, line_w=0.5)
    rect(s, cx, ry, 0.5, 0.92, TEAL)
    txt(s, num, cx + 0.01, ry + 0.26, 0.48, 0.38,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, cx + 0.6, ry + 0.05, 3.65, 0.32, size=12, bold=True, color=NAVY)
    txt(s, desc,  cx + 0.6, ry + 0.40, 3.65, 0.44, size=9.5, color=GRAY, wrap=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — 프로그램 소개 (캡처 공간 포함)
# 레이아웃: 왼쪽 기능 카드 4개 / 오른쪽 전체 화면 캡처
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 10, 5.625, WHITE)
header(s, "01", "프로그램 소개", "MSDS PDF → GHS 경고표지 자동 생성 도구")

# 왼쪽: 소개 + 기능 4개
txt(s, "MSDS PDF에서 경고표지 필수 항목을 자동 추출하여\n인쇄 가능한 라벨을 즉시 생성합니다.",
    0.35, 1.02, 4.5, 0.6, size=11, color=DARK, wrap=True)

features = [
    ("⚡", "자동 추출",  "PDF 텍스트에서 제품명·신호어\n유해문구 자동 인식"),
    ("✏",  "수동 편집", "추출 결과 직접 수정\n정확도 보완 가능"),
    ("🖨",  "즉시 출력", "HTML 생성 →\n브라우저 바로 인쇄"),
    ("📐", "규격 선택",  "아이라벨 3종 규격\n자동 레이아웃 적용"),
]
for i, (icon, title, desc) in enumerate(features):
    fx = 0.35 + (i % 2) * 2.2
    fy = 1.72 + (i // 2) * 1.55
    rect(s, fx, fy, 2.05, 1.38, LIGHT, line_rgb=LGRAY)
    rect(s, fx, fy, 2.05, 0.06, TEAL2)
    txt(s, icon, fx, fy + 0.1, 2.05, 0.45, size=24, align=PP_ALIGN.CENTER, color=TEAL)
    txt(s, title, fx, fy + 0.56, 2.05, 0.3, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txt(s, desc, fx + 0.08, fy + 0.86, 1.9, 0.45, size=9.5, color=GRAY, align=PP_ALIGN.CENTER, wrap=True)

# 오른쪽: 전체 화면 캡처 (메인 UI)
screenshot_box(s, 4.72, 1.02, 5.0, 4.3, "프로그램 전체 화면")

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — 화면 구성 (캡처 공간 + 3패널 설명)
# 레이아웃: 상단 전체 캡처 / 하단 3패널 설명
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 10, 5.625, LIGHT)
header(s, "02", "화면 구성", "3개 패널 레이아웃 — PDF 뷰어 / 편집 / 미리보기")

# 상단: 전체 화면 캡처 (넓게)
screenshot_box(s, 0.35, 1.05, 9.3, 2.72, "프로그램 전체 화면 캡처")

# 하단: 패널별 설명 3개
panels = [
    ("📄", "PDF 뷰어",        NAVY,                    ["PDF 파일 불러오기", "페이지 이동 (◀ ▶)", "원문 텍스트 보기"]),
    ("📋", "추출 결과 편집",  RGBColor(0x14,0x5A,0x72), ["제품명·신호어·H·P코드", "GHS 그림문자 선택", "글자 크기 조절"]),
    ("🏷", "경고표지 미리보기", RGBColor(0x0D,0x5C,0x4A), ["실시간 라벨 미리보기", "인쇄 HTML 생성 버튼", "브라우저 인쇄 연동"]),
]
pw = 2.95
for i, (icon, title, hbg, items) in enumerate(panels):
    px = 0.35 + i * 3.1
    rect(s, px, 3.95, pw, 1.4, WHITE, line_rgb=LGRAY)
    rect(s, px, 3.95, pw, 0.36, hbg)
    txt(s, f"{icon}  {title}", px + 0.12, 4.02, pw - 0.18, 0.24,
        size=10.5, bold=True, color=WHITE)
    for j, item in enumerate(items):
        txt(s, f"▸  {item}", px + 0.15, 4.4 + j * 0.32, pw - 0.22, 0.3,
            size=10, color=DARK, wrap=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — Step 1: PDF 불러오기 (캡처 공간 포함)
# 레이아웃: 왼쪽 4단계 / 오른쪽 PDF 뷰어 캡처
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 10, 5.625, WHITE)
header(s, "03", "Step 1 — PDF 불러오기", "MSDS PDF 파일 선택 및 텍스트 자동 추출")

steps = [
    ("PDF 파일 선택",    "[PDF 불러오기] 버튼 클릭\n파일 탐색기에서 MSDS PDF 선택"),
    ("자동 텍스트 추출", "PyMuPDF로 전체 텍스트 파싱\n제품명·신호어·H·P코드 자동 인식"),
    ("PDF 뷰어 확인",   "좌측 패널에 PDF 이미지 표시\n◀ ▶ 버튼으로 페이지 탐색"),
    ("원문 텍스트 확인", "[원문 보기] 버튼 → 추출 원문 확인\n인식 오류 파악에 활용"),
]
for i, (title, desc) in enumerate(steps):
    ry = 1.08 + i * 1.08
    rect(s, 0.35, ry, 4.65, 0.92, LIGHT, line_rgb=LGRAY, line_w=0.5)
    num_circle(s, i + 1, 0.42, ry + 0.22)
    txt(s, title, 0.95, ry + 0.06, 3.95, 0.3, size=12, bold=True, color=NAVY)
    txt(s, desc,  0.95, ry + 0.38, 3.95, 0.5, size=10, color=GRAY, wrap=True)

# 오른쪽: PDF 뷰어 패널 캡처
screenshot_box(s, 5.2, 1.05, 4.45, 4.25, "PDF 뷰어 패널 캡처")

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — Step 2: 추출 결과 편집 (캡처 공간 포함)
# 레이아웃: 왼쪽 항목 설명 / 오른쪽 편집 패널 캡처
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 10, 5.625, WHITE)
header(s, "04", "Step 2 — 추출 결과 편집", "자동 추출된 항목 확인 및 직접 수정")

fields = [
    ("제품명",        "추출된 화학물질명.\n오인식 시 직접 수정",                 NAVY),
    ("신호어",        "위험(빨강) / 경고(주황) /\n없음 — GHS 기준 자동 분류",    RED),
    ("유해·위험문구", "H코드 기반 자동 매칭\n(예: H225 인화성 액체)",            RGBColor(0x14,0x5A,0x72)),
    ("예방조치문구",  "P코드 기반 자동 매칭\n(예: P210 열원 멀리할 것)",         RGBColor(0x14,0x5A,0x72)),
    ("공급자 정보",   "제조사명·주소·긴급연락처.\n자유 편집 가능",                GRAY),
    ("GHS 그림문자",  "9종 중 해당 항목 체크.\n글자 크기도 개별 조절 가능",      TEAL),
]
for i, (name, desc, color) in enumerate(fields):
    fx = 0.35 + (i % 2) * 2.3
    fy = 1.05 + (i // 2) * 1.46
    rect(s, fx, fy, 2.15, 1.3, LIGHT, line_rgb=LGRAY)
    rect(s, fx, fy, 2.15, 0.06, color)
    txt(s, name, fx + 0.1, fy + 0.12, 2.0, 0.28, size=11.5, bold=True, color=color)
    txt(s, desc, fx + 0.1, fy + 0.46, 2.0, 0.78, size=10, color=GRAY, wrap=True)

# 오른쪽: 편집 패널 캡처
screenshot_box(s, 4.95, 1.05, 4.7, 4.28, "편집 패널 캡처")

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — Step 3: 인쇄 출력 (캡처 공간 포함)
# 레이아웃: 왼쪽 미리보기 캡처 / 오른쪽 출력 절차 4단계
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 10, 5.625, WHITE)
header(s, "05", "Step 3 — 인쇄 출력", "경고표지 미리보기 확인 → 인쇄용 HTML 생성 → 브라우저 인쇄")

# 왼쪽: 미리보기 패널 캡처
screenshot_box(s, 0.35, 1.05, 4.45, 4.28, "경고표지 미리보기 패널 캡처")

# 오른쪽: 출력 단계
print_steps = [
    ("미리보기 최종 확인", "우측 패널에서 제품명·신호어·\n그림문자·문구 내용 점검"),
    ("HTML 생성 버튼 클릭", "🖨 [인쇄 HTML 생성] 클릭\nprint_label.html 자동 생성"),
    ("브라우저에서 인쇄",  "Edge/Chrome 자동 실행 → Ctrl+P\n여백 '없음', 배율 100% 설정"),
    ("라벨지에 출력",      "아이라벨 전용지 트레이 삽입\n규격(칸 수) 확인 후 인쇄"),
]
for i, (title, desc) in enumerate(print_steps):
    ry = 1.08 + i * 1.08
    rect(s, 5.2, ry, 4.45, 0.92, LIGHT, line_rgb=LGRAY, line_w=0.5)
    num_circle(s, i + 1, 5.27, ry + 0.22)
    txt(s, title, 5.78, ry + 0.06, 3.75, 0.3, size=12, bold=True, color=NAVY)
    txt(s, desc,  5.78, ry + 0.38, 3.75, 0.5, size=10, color=GRAY, wrap=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — 라벨 규격 선택
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 10, 5.625, LIGHT)
header(s, "06", "라벨 규격 선택", "아이라벨 3종 규격 — 용도에 맞게 선택")

specs = [
    {"name": "CL233MP", "label": "9칸",  "w": "62mm",  "h": "91mm",  "grid": "3×3", "use": "소용량 시약병\n실험실 소분 용기",  "bg": TEAL},
    {"name": "CL835MP", "label": "24칸", "w": "60mm",  "h": "35mm",  "grid": "3×8", "use": "소형 샘플·바이알\n시험관 라벨",    "bg": RGBColor(0x14,0x5A,0x72)},
    {"name": "CL812MP", "label": "2칸",  "w": "196mm", "h": "125mm", "grid": "1×2", "use": "대형 드럼·탱크\n원료 보관 용기",   "bg": NAVY},
]
for i, sp in enumerate(specs):
    cx = 0.4 + i * 3.1
    rect(s, cx, 1.05, 2.9, 4.22, WHITE, line_rgb=LGRAY)
    rect(s, cx, 1.05, 2.9, 0.52, sp["bg"])
    txt(s, f"아이라벨  {sp['name']}", cx + 0.12, 1.1, 2.65, 0.26, size=12, bold=True, color=WHITE)
    txt(s, f"({sp['label']})", cx + 0.12, 1.37, 2.65, 0.18, size=10, color=RGBColor(0xAD,0xD8,0xE6))
    # 라벨 크기 시각화
    vis_w = float(sp["w"].replace("mm","")) / 220 * 2.3
    vis_h = float(sp["h"].replace("mm","")) / 140 * 1.0
    vis_x = cx + (2.9 - vis_w) / 2
    rect(s, vis_x, 1.72, vis_w, vis_h, LIGHT, line_rgb=sp["bg"], line_w=1.0)
    txt(s, f"{sp['w']} × {sp['h']}", cx + 0.1, 2.85, 2.7, 0.28,
        size=11, bold=True, color=sp["bg"], align=PP_ALIGN.CENTER)
    txt(s, f"배열: {sp['grid']}", cx + 0.1, 3.15, 2.7, 0.26,
        size=10.5, color=GRAY, align=PP_ALIGN.CENTER)
    rect(s, cx + 0.15, 3.55, 2.6, 0.03, LGRAY)
    txt(s, "권장 용도", cx + 0.15, 3.67, 2.6, 0.24, size=10, bold=True, color=GRAY)
    txt(s, sp["use"], cx + 0.15, 3.93, 2.6, 0.55, size=10.5, color=DARK, wrap=True)

    # 하단 캡처 공간 (해당 규격 실제 라벨 출력 캡처)
    screenshot_box(s, cx + 0.15, 4.55, 2.6, 0.62, f"{sp['label']} 출력 예시")

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — GHS 그림문자
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 10, 5.625, WHITE)
header(s, "07", "GHS 그림문자", "9종 — 해당 항목에 체크하여 경고표지에 포함")

ghs_items = [
    ("💣", "폭발성",    "GHS01", "자기반응성·유기과산화물"),
    ("🔥", "인화성",    "GHS02", "인화성 가스·액체·고체"),
    ("🔵", "산화성",    "GHS03", "산화성 물질"),
    ("🫙", "고압가스",  "GHS04", "압축·용해·냉동 액화가스"),
    ("⚗",  "부식성",   "GHS05", "금속·피부·눈 부식"),
    ("☠",  "급성독성", "GHS06", "경구·경피·흡입 독성"),
    ("🌿", "환경유해",  "GHS09", "수생환경 유해"),
    ("❗", "건강유해",  "GHS08", "발암성·생식독성"),
    ("⚠",  "경고",     "GHS07", "피부·눈 자극\n호흡기 민감성"),
]
cols_n = 5
for i, (icon, name, code, desc) in enumerate(ghs_items):
    gx = 0.3 + (i % cols_n) * 1.88
    gy = 1.05 + (i // cols_n) * 2.0
    rect(s, gx, gy, 1.72, 1.72, LIGHT, line_rgb=LGRAY)
    rect(s, gx, gy, 1.72, 0.06, TEAL2)
    txt(s, icon, gx, gy + 0.1, 1.72, 0.48, size=26, align=PP_ALIGN.CENTER, color=TEAL)
    txt(s, name, gx, gy + 0.62, 1.72, 0.26, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txt(s, code, gx, gy + 0.9, 1.72, 0.22, size=9, color=TEAL, align=PP_ALIGN.CENTER)
    txt(s, desc, gx + 0.06, gy + 1.13, 1.6, 0.5, size=8.5, color=GRAY, align=PP_ALIGN.CENTER, wrap=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — 주의사항 & FAQ
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 10, 5.625, WHITE)
header(s, "08", "주의사항 & FAQ", "올바른 사용을 위한 체크리스트")

# 왼쪽: 주의사항
rect(s, 0.35, 1.05, 4.55, 4.25, LIGHT, line_rgb=LGRAY)
rect(s, 0.35, 1.05, 4.55, 0.4, RED)
txt(s, "⚠  주의사항", 0.5, 1.12, 4.2, 0.26, size=12, bold=True, color=WHITE)
cautions = [
    "자동 추출 결과는 참고용입니다.\n최종 인쇄 전 원본 MSDS와 반드시 대조하세요.",
    "스캔 PDF는 인식률이 낮습니다.\n텍스트 레이어가 있는 원본 PDF를 사용하세요.",
    "법적 효력이 있는 라벨은 안전보건공단\n규정에 따라 별도 검토가 필요합니다.",
    "인쇄 전 라벨 규격(아이라벨 종류)이\n실제 라벨지와 일치하는지 확인하세요.",
]
for j, c in enumerate(cautions):
    rect(s, 0.45, 1.57 + j * 0.88, 4.35, 0.78, WHITE, line_rgb=RGBColor(0xF5,0xC6,0xCB), line_w=0.5)
    txt(s, "●", 0.52, 1.66 + j * 0.88, 0.22, 0.28, size=10, color=RED)
    txt(s, c, 0.75, 1.63 + j * 0.88, 3.95, 0.65, size=10, color=DARK, wrap=True)

# 오른쪽: FAQ
rect(s, 5.1, 1.05, 4.55, 4.25, LIGHT, line_rgb=LGRAY)
rect(s, 5.1, 1.05, 4.55, 0.4, TEAL)
txt(s, "💬  자주 묻는 질문 (FAQ)", 5.22, 1.12, 4.2, 0.26, size=12, bold=True, color=WHITE)
faqs = [
    ("Q. 제품명이 잘못 추출됐어요",     "편집 패널의 제품명 칸을 직접 수정하면 됩니다."),
    ("Q. 그림문자가 자동 선택 안 돼요", "H코드 기반 자동 선택 후 미체크 항목은\n수동으로 체크하세요."),
    ("Q. 인쇄 시 라벨 위치가 안 맞아요","브라우저 인쇄 설정에서 여백을\n'없음' 또는 '최소'로 설정하세요."),
    ("Q. 프로그램 실행이 안 돼요",      "MSDS_Label_Maker.exe 더블클릭 후\n브라우저가 자동으로 열립니다."),
]
for j, (q, a) in enumerate(faqs):
    ry = 1.57 + j * 0.92
    rect(s, 5.18, ry, 4.38, 0.82, WHITE, line_rgb=LGRAY, line_w=0.5)
    txt(s, q, 5.28, ry + 0.05, 4.18, 0.28, size=10.5, bold=True, color=TEAL)
    txt(s, a, 5.28, ry + 0.36, 4.18, 0.42, size=10, color=DARK, wrap=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — 마무리
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 10, 5.625, NAVY)
rect(s, 0, 0, 0.18, 5.625, TEAL)
rect(s, 9.82, 0, 0.18, 5.625, TEAL)

txt(s, "🧪", 4.55, 1.2, 1.0, 0.85, size=40, align=PP_ALIGN.CENTER, color=TEAL)
txt(s, "MSDS Label Maker v2.1.0",
    1.0, 2.1, 8.0, 0.58, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "GHS 경고표지 자동 생성 도구",
    1.0, 2.7, 8.0, 0.38, size=16, color=TEAL2, align=PP_ALIGN.CENTER)
rect(s, 3.8, 3.2, 2.4, 0.04, TEAL2)
txt(s, "LX글라스㈜  연구기획팀  |  박재영  |  2026-06-11",
    1.0, 3.35, 8.0, 0.32, size=12, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "문의 및 개선 요청은 연구기획팀으로 연락해 주세요.",
    1.0, 3.76, 8.0, 0.3, size=11, color=GRAY, align=PP_ALIGN.CENTER)
rect(s, 2.5, 4.55, 5.0, 0.72, TEAL)
txt(s, "사용 중 오류 발생 시  →  app_error.txt 내용과 함께 문의",
    2.6, 4.67, 4.8, 0.34, size=10.5, color=WHITE, align=PP_ALIGN.CENTER)

out = r"C:\-msds_label_maker\MSDS_Label_Maker_사용설명서_v2.1.0.pptx"
prs.save(out)
print(f"저장 완료: {out}")
